import json
import logging
import os
import shutil
import subprocess
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timezone
from typing import Any
from .base_agent import AbstractStudyAgent, TOOL_DEFINITIONS

log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Phase 1 — NotesAgent
# ---------------------------------------------------------------------------

class NotesAgent(AbstractStudyAgent):
    """Generates structured Markdown notes (notes.md) + timing.json sidecar.

    Equipped with web_search and image_search tools so the model can verify
    facts and embed relevant diagram URLs directly in the output.
    """

    def __init__(
        self,
        model: str = "claude-sonnet-4-6",
        api_key: str | None = None,
    ) -> None:
        super().__init__(model=model, api_key=api_key)
        self.tools.append(TOOL_DEFINITIONS["web_search"])
        self.tools.append(TOOL_DEFINITIONS["image_search"])

    def build_prompt(self, topic: str) -> str:
        """Return the notes-generation prompt for the given topic.

        Args:
            topic: The subject to write notes about.

        Returns:
            A fully formatted prompt string.
        """
        return (
            f'Write comprehensive study notes on: "{topic}".\n\n'
            "INSTRUCTIONS — follow all steps in order:\n\n"
            "1. WEB SEARCH (strongly preferred): Call web_search at least twice "
            "before writing. If the first query returns no results, retry with a "
            "broader query. Ground every major concept in something you found.\n\n"
            "2. IMAGE EMBEDDING: For every ## section heading, call image_search "
            "and embed the best result immediately below the heading:\n"
            "   ![description](url)\n"
            "   *Caption: one-line description.*\n"
            "If image_search fails or returns no results, skip the image for that "
            "section and continue — do NOT stop generating notes.\n\n"
            "3. FALLBACK: If ALL web_search calls fail or return errors, write the "
            "notes from your internal knowledge. Still produce complete notes — "
            "do not output any error messages or refuse to generate content.\n\n"
            "FORMAT REQUIREMENTS:\n"
            "- Use ## headers for each major concept (required)\n"
            "- Bullet-point key definitions and properties under each header\n"
            "- 1-2 concrete examples per section, labelled **Example:**\n"
            "- 400-700 words total, Markdown only, no preamble before the first header\n\n"
            "AFTER your Markdown notes, on its own line write exactly:\n"
            "---TIMING---\n"
            "Then write a JSON array, one object per ## section:\n"
            '[{"section": "heading without ##", "narration": "1-2 sentence summary", '
            '"estimated_seconds": <int 20-60>}]\n'
            "Return ONLY the JSON array after the separator — no code fences, no explanation.\n"
        )

    def run(self, output_dir: str | None = None, **kwargs: Any) -> dict:
        """Generate notes, save notes.md + timing.json sidecar.

        Keyword Args:
            topic (str):        The subject to generate notes about.
            output_dir (str):   Absolute path to the run's output directory.
                                Both output files are written flat into this
                                directory as ``notes.md`` and ``timing.json``.
                                Falls back to ``output/`` at the project root
                                when None.

        Returns:
            ``{"status": "ok", "output": <notes content>, "md_path": str,
               "timing_path": str}``
        """
        topic: str = kwargs["topic"]

        if output_dir is None:
            _base = os.path.dirname(
                os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            )
            output_dir = os.path.join(_base, "output")
        os.makedirs(output_dir, exist_ok=True)

        t0 = time.monotonic()
        started_ts = datetime.now(timezone.utc).isoformat()

        content = self._call_api(self.build_prompt(topic), use_tools=True)

        duration = round(time.monotonic() - t0, 3)
        finished_ts = datetime.now(timezone.utc).isoformat()

        result: dict = {"status": "ok", "output": content, "md_path": "", "timing_path": ""}

        if not self.validate_output(result):
            log.error(
                "NotesAgent validate_output FAILED.\n"
                "  has '## ': %s\n"
                "  word count: %d\n"
                "  first 300 chars: %r",
                "## " in content,
                len(content.split()),
                content[:300],
            )
            return {"status": "error", "output": content, "md_path": "", "timing_path": ""}

        # Attempt to extract timing data from the combined output (saves a second LLM call).
        # build_prompt instructs Claude to append ---TIMING--- then a JSON array.
        sections: list = []
        if "---TIMING---" in content:
            notes_part, timing_raw = content.split("---TIMING---", 1)
            content = notes_part.strip()
            timing_raw = (
                timing_raw.strip()
                .lstrip("```json")
                .lstrip("```")
                .rstrip("```")
                .strip()
            )
            try:
                sections = json.loads(timing_raw)
                log.info("[Notes] Timing extracted inline — skipping second LLM call.")
            except json.JSONDecodeError:
                log.warning("[Notes] Inline timing JSON malformed — falling back to second call.")

        if not sections:
            # Fallback: dedicated second LLM call (original behaviour).
            # VideoAgent._build_narration_scripts and _generate_html_slides both
            # iterate over this list and key into "section", "narration",
            # "estimated_seconds" — the flat metadata dict would cause a KeyError.
            timing_prompt = (
                "Given these study notes, return a JSON array where each element "
                "represents one ## section. Each object must have exactly these keys:\n"
                '  "section": the heading text (without the ## prefix),\n'
                '  "narration": 1-2 sentence summary of what to say about this section,\n'
                '  "estimated_seconds": integer seconds to speak about it '
                "(typically 20-60 per section).\n\n"
                "Return ONLY valid JSON — no markdown fences, no explanation.\n\n"
                f"NOTES:\n{content}"
            )
            timing_raw = self._call_api(timing_prompt, use_tools=False).strip()
            if timing_raw.startswith("```"):
                lines = timing_raw.split("\n", 1)
                timing_raw = lines[1].rsplit("```", 1)[0].strip() if len(lines) > 1 else ""
            try:
                sections = json.loads(timing_raw)
            except json.JSONDecodeError:
                sections = []

        md_path = self._save_output(content, "notes.md", output_dir=output_dir)
        timing_path = self._save_output(
            json.dumps(
                {
                    "agent_id": self.agent_id,
                    "topic": topic,
                    "started_at": started_ts,
                    "finished_at": finished_ts,
                    "duration_seconds": duration,
                    "sections": sections,
                },
                indent=2,
            ),
            "timing.json",
            output_dir=output_dir,
        )

        result["md_path"] = md_path
        result["timing_path"] = timing_path
        result["duration_s"] = duration
        result["started_at"] = started_ts
        result["finished_at"] = finished_ts
        print(f"[Notes] Saved:  {md_path}")
        print(f"[Notes] Timing: {timing_path}")
        return result

    def validate_output(self, output: dict) -> bool:
        """Return True if the notes content passes structural and length checks.

        Args:
            output: The dict from run() — inspects output["output"].
        """
        if output.get("status") == "error":
            return False
        content: str = output.get("output", "")
        if not isinstance(content, str) or not content.strip():
            return False
        if "## " not in content:
            return False
        return 200 <= len(content.split()) <= 2000


# ---------------------------------------------------------------------------
# Phase 2a — FlashcardAgent
# ---------------------------------------------------------------------------

class FlashcardAgent(AbstractStudyAgent):
    """Transforms notes.md into Obsidian-compatible spaced-repetition flashcards.

    No tools needed — the agent reads the provided notes and rephrases them
    into structured active-recall cards without consulting external sources.
    """

    def __init__(
        self,
        model: str = "claude-haiku-4-5-20251001",
        api_key: str | None = None,
    ) -> None:
        super().__init__(model=model, api_key=api_key)

    def build_prompt(self, notes_content: str) -> str:
        """Return the flashcard-generation prompt.

        Instructs the model to produce 8-12 Obsidian-style flashcards with a
        prescribed mix of definition, application, and distinction cards.

        Args:
            notes_content: Full text of notes.md from Phase 1.

        Returns:
            A fully formatted prompt string.
        """
        return (
            "Read the following study notes carefully, then generate between "
            "12 and 15 flashcards.\n\n"
            "NOTES:\n"
            f"{notes_content}\n\n"
            "Use this exact Markdown format for every card — no deviations:\n\n"
            "## {{question}} #flashcard\n\n"
            "{{answer — see ANSWER REQUIREMENTS below}}\n\n"
            "---\n\n"
            "ANSWER REQUIREMENTS — all four rules apply to every single card:\n"
            "  1. Minimum 3 sentences per answer — no exceptions.\n"
            "  2. Every answer must include EITHER a concrete real-world example "
            "OR an explicit comparison to a related concept.\n"
            "  3. End every answer with a sentence starting exactly 'Why it matters:' "
            "that explains the practical relevance of the concept.\n"
            "  4. Every answer must be derivable directly from the notes above "
            "— introduce no new information.\n\n"
            "Card type mix (all three types required, totalling 12-15):\n"
            "  • 4-5 definition cards    — question form: 'What is X?'\n"
            "  • 4-5 application cards   — question form: "
            "'How does X work in the context of Y?'\n"
            "  • 3-5 distinction cards   — question form: "
            "'What is the difference between X and Y?'\n\n"
            "Additional constraints:\n"
            "  • Do not copy sentences verbatim from the notes — rephrase in "
            "your own words.\n"
            "  • Return ONLY the Markdown flashcard content — no preamble, "
            "no closing remarks, nothing else.\n"
        )

    def run(self, notes_content: str, output_dir: str = "output") -> dict:
        """Generate flashcards from notes and save them to flashcards.md.

        Args:
            notes_content: Full text of notes.md from Phase 1.
            output_dir:    Subdirectory prefix passed to _save_output.

        Returns:
            ``{"status": "ok", "flashcards_path": str, "flashcards_content": str}``
        """
        t0 = time.monotonic()
        started_at = datetime.now(timezone.utc).isoformat()
        prompt = self.build_prompt(notes_content)
        response = self._call_api(prompt, use_tools=False)
        duration_s = round(time.monotonic() - t0, 3)
        finished_at = datetime.now(timezone.utc).isoformat()
        # Write into the run directory as flashcards.md — api_server.py maps
        # "flashcards_md" -> "flashcards.md" and serves it from run_dir, so a
        # global output/flashcards/<uuid>.md path is unreachable by the frontend.
        flashcards_path = self._save_output(
            response, "flashcards.md", output_dir=output_dir
        )
        print(f"[Flashcards] Saved: {flashcards_path}")
        return {
            "status": "ok",
            "flashcards_path": flashcards_path,
            "flashcards_content": response,
            "duration_s": duration_s,
            "started_at": started_at,
            "finished_at": finished_at,
        }

    def validate_output(self, output: dict) -> bool:
        """Return True if the flashcard content meets structural requirements.

        Passes only when the content contains at least 6 H2 headers, 6
        #flashcard tags, and 6 horizontal-rule separators — a proxy for
        having produced at least 6 well-formed cards.

        Args:
            output: The dict from run() — inspects output["flashcards_content"].
        """
        content: str = output.get("flashcards_content", "")
        if not isinstance(content, str) or not content.strip():
            return False
        return (
            content.count("## ") >= 6
            and content.count("#flashcard") >= 6
            and content.count("---") >= 6
        )


# ---------------------------------------------------------------------------
# Phase 2b — VideoAgent
# ---------------------------------------------------------------------------

class VideoAgent(AbstractStudyAgent):
    """Converts notes.md + section timing data into a narrated MP4 study video.

    Pipeline:
      1. _build_narration_scripts — expand timing entries into full narrations (LLM)
      2. _generate_html_slides    — LLM-generated HTML slides → Playwright PNG screenshots
      3. _generate_audio          — OpenAI TTS, one MP3 per narration
      4. _assemble_video          — MoviePy stitches frames + audio into MP4
    """

    def __init__(
        self,
        model: str = "claude-haiku-4-5-20251001",
        api_key: str | None = None,
        openai_api_key: str | None = None,
    ) -> None:
        """Initialise the agent.

        Args:
            model:          Anthropic model ID.
            api_key:        Anthropic API key.
            openai_api_key: OpenAI API key for TTS (falls back to OPENAI_API_KEY env var).
        """
        super().__init__(model=model, api_key=api_key)
        self.openai_api_key = openai_api_key

    def build_prompt(self, topic: str) -> str:
        """No-op — VideoAgent builds all prompts inline in its private methods.

        Args:
            topic: Unused; present to satisfy the abstract interface.

        Returns:
            An empty string.
        """
        return ""

    def run(
        self,
        notes_content: str,
        timing_json: list,
        output_dir: str | None = None,
    ) -> dict:
        """Generate a narrated study video and PPTX from notes and timing data.

        Args:
            notes_content: Full text of notes.md from Phase 1.
            timing_json:   List of section dicts, each with keys:
                           "section", "narration", "estimated_seconds".
            output_dir:    Absolute path to the run's output directory.
                           ``study_video.mp4`` is written directly here;
                           slides go into ``output_dir/slides/`` and audio
                           into ``output_dir/audio/``.
                           Falls back to ``output/`` at the project root when
                           None.

        Returns:
            ``{"status": "ok", "video_path": str, "pptx_path": str}``

        Raises:
            AssertionError: If frame count and narration count diverge after
                            slide generation.
        """
        if output_dir is None:
            _base = os.path.dirname(
                os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            )
            output_dir = os.path.join(_base, "output")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, "study_video.mp4")

        _t0 = time.monotonic()

        # Stage A: narrations and slides both need only timing_json — run in parallel
        t = time.monotonic()
        with ThreadPoolExecutor(max_workers=2) as _stage_a:
            _narr_fut  = _stage_a.submit(self._build_narration_scripts, notes_content, timing_json)
            _slide_fut = _stage_a.submit(self._generate_html_slides, timing_json, notes_content, output_dir)
            narrations  = _narr_fut.result()
            frame_paths = _slide_fut.result()
        dt_narrations = time.monotonic() - t  # combined wall-clock for stage A
        dt_slides = dt_narrations              # same wall-clock (ran in parallel)

        assert len(frame_paths) == len(narrations), (
            f"Frame/narration count mismatch: {len(frame_paths)} frames vs "
            f"{len(narrations)} narrations."
        )

        # Stage B+C: audio needs narrations; PPTX needs frame_paths — both ready, run in parallel
        t = time.monotonic()
        with ThreadPoolExecutor(max_workers=2) as _stage_bc:
            _audio_fut = _stage_bc.submit(self._generate_audio, narrations, output_dir)
            _pptx_fut  = _stage_bc.submit(self._export_pptx, frame_paths, timing_json, output_dir)
            audio_paths = _audio_fut.result()
            pptx_path   = _pptx_fut.result()
        dt_audio = time.monotonic() - t  # combined wall-clock for stage B+C
        dt_pptx  = dt_audio              # same wall-clock (ran in parallel)

        # Stage D: video assembly needs both frame_paths and audio_paths
        t = time.monotonic()
        final_path = self._assemble_video(frame_paths, audio_paths, output_path)
        dt_assemble = time.monotonic() - t

        dt_total = time.monotonic() - _t0
        print(
            f"\n[Video] Benchmark ({len(timing_json)} sections):\n"
            f"  stage-A (narrations‖slides) : {dt_narrations:6.1f}s\n"
            f"  stage-BC (audio‖pptx)       : {dt_audio:6.1f}s\n"
            f"  assemble                    : {dt_assemble:6.1f}s\n"
            f"  total                       : {dt_total:6.1f}s"
        )

        result: dict = {
            "status": "ok",
            "video_path": final_path,
            "pptx_path": pptx_path,
            "benchmark": {
                "stage_a_s": round(dt_narrations, 2),   # narrations ‖ slides (wall-clock)
                "stage_bc_s": round(dt_audio, 2),        # audio ‖ pptx (wall-clock)
                "assemble_s": round(dt_assemble, 2),
                "total_s": round(dt_total, 2),
            },
        }
        if not self.validate_output(result):
            return {"status": "error", "video_path": final_path, "pptx_path": pptx_path}

        print(f"[Video] Saved to:      {final_path}")
        print(f"[Video] PPTX saved to: {pptx_path}")
        return result

    def validate_output(self, output: dict) -> bool:
        """Return True if both the MP4 and the PPTX exist and are non-empty.

        Args:
            output: The dict from run() — inspects ``output["video_path"]``
                    and ``output["pptx_path"]``.
        """
        video_path: str = output.get("video_path", "")
        pptx_path: str = output.get("pptx_path", "")
        video_ok = (
            isinstance(video_path, str)
            and video_path.endswith(".mp4")
            and os.path.isfile(video_path)
            and os.path.getsize(video_path) > 0
        )
        pptx_ok = (
            isinstance(pptx_path, str)
            and pptx_path.endswith(".pptx")
            and os.path.isfile(pptx_path)
            and os.path.getsize(pptx_path) > 0
        )
        return video_ok and pptx_ok

    # ------------------------------------------------------------------
    # Private helpers
    # ------------------------------------------------------------------

    def _build_narration_scripts(
        self, notes_content: str, timing_json: list
    ) -> list[str]:
        """Expand each timing entry into a full spoken narration via the LLM.

        All API calls are issued concurrently (one thread per section, capped
        at 5 workers to avoid rate-limit spikes).

        Args:
            notes_content: Full text of notes.md — context for every call.
            timing_json:   List of section dicts with keys "narration" and
                           "estimated_seconds".

        Returns:
            List of expanded narration strings, one per timing entry, in order.
        """
        def _narrate(i: int, entry: dict) -> tuple[int, str]:
            prompt = (
                f"Given this section from study notes: {entry['narration']}\n"
                f"Expand this into a spoken narration of approximately "
                f"{entry['estimated_seconds']} seconds.\n"
                f"Use the following full notes as context: {notes_content}\n"
                "Tone: clear, educational, like a university lecturer.\n"
                "Return ONLY the narration text, no labels or headers."
            )
            return i, self._call_api(prompt, use_tools=False)

        workers = min(len(timing_json), 5)
        narrations: list[str] = [""] * len(timing_json)
        with ThreadPoolExecutor(max_workers=workers) as pool:
            futures = {pool.submit(_narrate, i, entry): i for i, entry in enumerate(timing_json)}
            for fut in as_completed(futures):
                i, text = fut.result()
                narrations[i] = text
        return narrations

    def _generate_html_slides(
        self, timing_json: list, notes_content: str, output_dir: str
    ) -> list[str]:
        """Generate one HTML slide per timing entry and screenshot to PNG.

        LLM generation and Playwright screenshot for each slide are issued
        concurrently (one thread per slide, capped at 5 workers).

        Args:
            timing_json:   List of section dicts with keys "section" and
                           "narration".
            notes_content: Full notes text passed as LLM context.
            output_dir:    Absolute path to the run's output directory.
                           HTML and PNG files are placed in a ``slides/``
                           subdirectory within this directory.

        Returns:
            List of absolute PNG file paths, one per timing entry, in order.
        """
        def _make_slide(i: int, entry: dict) -> tuple[int, str]:
            prompt = (
                f"Generate a single self-contained HTML slide for the concept: "
                f"{entry['section']}\n"
                "The slide should:\n"
                "- Have a clean white background, 1280x720px\n"
                "- Show the section title in large text at top\n"
                "- Show 3-5 bullet points summarising the key ideas from the "
                f"narration: {entry['narration']}\n"
                "- Use simple inline CSS only — no external libraries\n"
                "- Be a complete HTML document that renders as a slide when "
                "opened in a browser\n"
                "Return ONLY the HTML, no explanation."
            )
            html_content = self._call_api(prompt, use_tools=False)
            html_path = self._save_output(
                html_content, f"slides/slide_{i:02d}.html", output_dir=output_dir
            )
            png_path = os.path.join(os.path.dirname(html_path), f"slide_{i:02d}.png")
            self._html_to_png(html_path, png_path)
            return i, png_path

        workers = min(len(timing_json), 5)
        png_paths: list[str] = [""] * len(timing_json)
        with ThreadPoolExecutor(max_workers=workers) as pool:
            futures = {pool.submit(_make_slide, i, entry): i for i, entry in enumerate(timing_json)}
            for fut in as_completed(futures):
                i, png_path = fut.result()
                png_paths[i] = png_path
        return png_paths

    def _html_to_png(self, html_path: str, output_png: str) -> str:
        """Convert an HTML file to a 1280×720 PNG using Playwright sync API.

        Uses the synchronous Playwright API so this method is safe to call
        from ThreadPoolExecutor threads (which have no running event loop)
        and from within ADK's async runner context (where asyncio.run would
        raise "This event loop is already running").

        Args:
            html_path:  Absolute path to the .html file.
            output_png: Destination path for the .png screenshot.

        Returns:
            output_png path after successful conversion.
        """
        from playwright.sync_api import sync_playwright

        with sync_playwright() as pw:
            browser = pw.chromium.launch()
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            page.goto(f"file://{html_path}")
            page.screenshot(path=output_png, full_page=False)
            browser.close()
        return output_png

    def _generate_audio(self, narrations: list[str], output_dir: str) -> list[str]:
        """Generate one TTS MP3 per narration using OpenAI tts-1-hd (voice: coral).

        All TTS requests are issued concurrently (one thread per narration,
        capped at 5 workers).

        Args:
            narrations: List of narration strings, one per slide.
            output_dir: Absolute path to the run's output directory.
                        MP3 files are placed in an ``audio/`` subdirectory
                        within this directory.

        Returns:
            List of MP3 file paths in the same order as narrations.
        """
        import openai

        audio_dir = os.path.join(output_dir, "audio")
        os.makedirs(audio_dir, exist_ok=True)

        client = openai.OpenAI(api_key=self.openai_api_key)

        def _tts(i: int, narration: str) -> tuple[int, str]:
            out_path = os.path.join(audio_dir, f"audio_{i:02d}.mp3")
            with client.audio.speech.with_streaming_response.create(
                model="tts-1-hd",
                voice="coral",
                input=narration,
            ) as response:
                response.stream_to_file(out_path)
            return i, out_path

        workers = min(len(narrations), 5)
        audio_paths: list[str] = [""] * len(narrations)
        with ThreadPoolExecutor(max_workers=workers) as pool:
            futures = {pool.submit(_tts, i, n): i for i, n in enumerate(narrations)}
            for fut in as_completed(futures):
                i, path = fut.result()
                audio_paths[i] = path
        return audio_paths

    def _export_pptx(
        self,
        frame_paths: list[str],
        timing_json: list,
        output_dir: str,
    ) -> str:
        """Build a widescreen PPTX where each slide is a PNG frame + title overlay.

        Each PNG frame is placed as a full-bleed background image. A dark navy
        text box is overlaid at the top of each slide with the section title in
        white 28 pt bold text.  The file is saved as
        ``{output_dir}/slides/flashcards.pptx``.

        Requires the ``python-pptx`` package (``pip install python-pptx``).

        Args:
            frame_paths: Absolute PNG file paths in slide order.
            timing_json: Section timing dicts; key ``"section"`` is used as
                         the slide title.
            output_dir:  Absolute path to the run's output directory.

        Returns:
            Absolute path to the saved ``flashcards.pptx``.
        """
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]  # truly blank layout

        for frame_path, entry in zip(frame_paths, timing_json):
            slide = prs.slides.add_slide(blank_layout)

            # Full-bleed PNG background
            slide.shapes.add_picture(
                frame_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

            # Title overlay — dark navy box at the top, white bold 28 pt text
            txBox = slide.shapes.add_textbox(
                left=0, top=0, width=prs.slide_width, height=Inches(1.2)
            )
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)  # dark navy
            tf = txBox.text_frame
            tf.word_wrap = True
            run = tf.paragraphs[0].add_run()
            run.text = entry.get("section", "")
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        slides_dir = os.path.join(output_dir, "slides")
        os.makedirs(slides_dir, exist_ok=True)
        pptx_path = os.path.join(slides_dir, "flashcards.pptx")
        prs.save(pptx_path)
        return pptx_path

    def _assemble_video(
        self,
        frame_paths: list[str],
        audio_paths: list[str],
        output_path: str,
    ) -> str:
        """Combine PNG frames and MP3 clips into a single MP4 via MoviePy.

        Each (frame, audio) pair becomes one clip whose on-screen duration
        matches the length of its audio clip.

        Args:
            frame_paths: PNG file paths, one per slide.
            audio_paths: MP3 file paths, one per slide.
            output_path: Destination MP4 path.

        Returns:
            output_path after successful export.
        """
        from moviepy import AudioFileClip, ImageClip, concatenate_videoclips

        clips = []
        for frame, audio in zip(frame_paths, audio_paths):
            audio_clip = AudioFileClip(audio)
            clips.append(
                ImageClip(frame)
                .with_duration(audio_clip.duration)
                .with_audio(audio_clip)
            )

        final = concatenate_videoclips(clips, method="compose")
        final.write_videofile(output_path, fps=24, codec="libx264", audio_codec="aac")
        return output_path


# ---------------------------------------------------------------------------
# Phase 3 — PDFAgent
# ---------------------------------------------------------------------------

# Directories where pandoc / LaTeX engines commonly land on Windows but which
# are not always reflected in the *current process'* PATH (e.g. when a terminal
# was launched before the installer updated the environment). We search these
# explicitly so a stale PATH can never break PDF generation again.
_EXTRA_TOOL_DIRS = [
    os.path.join(os.environ.get("LOCALAPPDATA", ""), "Pandoc"),
    os.path.join(os.environ.get("ProgramFiles", ""), "Pandoc"),
    os.path.join(os.environ.get("ProgramFiles(x86)", ""), "Pandoc"),
    os.path.join(os.environ.get("LOCALAPPDATA", ""), "Microsoft", "WinGet", "Links"),
    os.path.join(os.environ.get("USERPROFILE", ""), "scoop", "shims"),
    os.path.join(os.environ.get("LOCALAPPDATA", ""), "Tectonic"),
]

# Candidate PDF engines, in order of preference. The first one found wins.
_PDF_ENGINES = ["xelatex", "lualatex", "pdflatex", "tectonic", "wkhtmltopdf", "weasyprint"]


def _resolve_tool(name: str) -> str | None:
    """Locate an executable by name, searching PATH and known install dirs.

    Returns the absolute path to the executable, or None if not found.
    Works even when the running process inherited a stale PATH.
    """
    found = shutil.which(name)
    if found:
        return found
    for d in _EXTRA_TOOL_DIRS:
        if not d:
            continue
        candidate = shutil.which(name, path=d)
        if candidate:
            return candidate
    return None


class PDFAgent(AbstractStudyAgent):
    """Converts a Markdown file to PDF via pandoc + xelatex.

    Does not call the LLM; build_prompt is a no-op to satisfy the ABC.
    """

    def __init__(
        self,
        model: str = "claude-sonnet-4-6",
        api_key: str | None = None,
    ) -> None:
        super().__init__(model=model, api_key=api_key)

    def build_prompt(self, topic: str) -> str:
        """PDFAgent uses pandoc, not LLM prompts."""
        return ""

    def run(
        self,
        input_md_path: str,
        output_pdf_path: str | None = None,
    ) -> dict:
        """Render a Markdown file to PDF with pandoc + xelatex.

        Args:
            input_md_path:   Absolute path to the source .md file.
            output_pdf_path: Destination .pdf path. If None, derived by
                             replacing the .md extension with .pdf.

        Returns:
            ``{"status": "ok", "pdf_path": str}``

        Raises:
            RuntimeError: If pandoc is not installed or exits non-zero.
        """
        if output_pdf_path is None:
            output_pdf_path = os.path.splitext(input_md_path)[0] + ".pdf"

        pandoc = _resolve_tool("pandoc")
        if pandoc is None:
            raise RuntimeError(
                "pandoc could not be found on PATH or in any known install "
                "location. Install it from https://pandoc.org/installing.html "
                "(or `winget install JohnMacFarlane.Pandoc`)."
            )

        engine = next((e for e in _PDF_ENGINES if _resolve_tool(e)), None)
        if engine is None:
            raise RuntimeError(
                "pandoc is installed, but no PDF engine was found. Install one "
                "of: a LaTeX engine such as tectonic "
                "(`winget install TectonicTypesetting.Tectonic`), MiKTeX, or "
                "TeX Live (provides xelatex/pdflatex), or wkhtmltopdf "
                "(`winget install wkhtmltopdf.wkhtmltox`)."
            )

        # Eisvogel is a LaTeX template, so only probe for it when the resolved
        # engine is a LaTeX one. HTML engines (wkhtmltopdf/weasyprint) ignore or
        # error on --template and on the -V geometry flags.
        _is_latex = engine in ("xelatex", "lualatex", "pdflatex", "tectonic")
        _eisvogel = None
        if _is_latex:
            _eisvogel_candidates = [
                os.path.expandvars(r"%APPDATA%\pandoc\templates\eisvogel.latex"),
                os.path.expandvars(r"%USERPROFILE%\.pandoc\templates\eisvogel.latex"),
                os.path.join(
                    os.path.expanduser("~"), ".pandoc", "templates", "eisvogel.latex"
                ),
            ]
            _eisvogel = next(
                (p for p in _eisvogel_candidates if os.path.isfile(p)), None
            )

        cmd = [
            pandoc, input_md_path,
            "-o", output_pdf_path,
            f"--pdf-engine={_resolve_tool(engine)}",
        ]
        if _eisvogel:
            cmd += [
                f"--template={_eisvogel}",
                "-V", "geometry:margin=0.75in",
                "-V", "linkcolor=blue",
                "-V", "urlcolor=blue",
                "-V", "colorlinks=true",
                "-V", "numbersections",
                "-V", "toc",
                "--highlight-style=tango",
                "--dpi=300",
            ]
            log.info("[PDF] Using Eisvogel template: %s", _eisvogel)
        elif _is_latex:
            cmd += [
                "-V", "geometry:margin=1in",
                "-V", "colorlinks=true",
                "-V", "linkcolor=blue",
                "--highlight-style=tango",
            ]
            log.info(
                "[PDF] Eisvogel template not found - using basic formatting. "
                "See the Eisvogel section in README.md to install it."
            )

        t0 = time.monotonic()
        started_at = datetime.now(timezone.utc).isoformat()
        try:
            # Decode output as UTF-8 with replacement: pandoc/LaTeX engines can
            # emit non-cp1252 bytes (e.g. tectonic's font diagnostics), which
            # would otherwise crash the capture thread on Windows.
            subprocess.run(
                cmd,
                check=True,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
            )
        except subprocess.CalledProcessError as exc:
            raise RuntimeError(
                f"pandoc failed (engine={engine}):\n{exc.stderr}"
            ) from exc

        duration_s = round(time.monotonic() - t0, 3)
        finished_at = datetime.now(timezone.utc).isoformat()
        result: dict = {
            "status": "ok",
            "pdf_path": output_pdf_path,
            "duration_s": duration_s,
            "started_at": started_at,
            "finished_at": finished_at,
        }
        if not self.validate_output(result):
            return {"status": "error", "pdf_path": output_pdf_path, "duration_s": duration_s}

        print(f"[PDF] Saved to: {output_pdf_path}")
        return result

    def validate_output(self, output: dict) -> bool:
        """Return True if the PDF exists on disk and has non-zero size.

        Args:
            output: The dict from run() — inspects output["pdf_path"].
        """
        path: str = output.get("pdf_path", "")
        return (
            isinstance(path, str)
            and path.endswith(".pdf")
            and os.path.isfile(path)
            and os.path.getsize(path) > 0
        )
