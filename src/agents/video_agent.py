"""
video_agent.py

Phase 2 of the pipeline: turns notes into a narrated slide video.

Split out of specialist_agent.py because this is the whole heavy-dependency
surface -- Playwright and Chromium for slide rendering, MoviePy and ffmpeg for
assembly, python-pptx for the deck, and OpenAI for text-to-speech. Keeping it
here means the notes, flashcard and PDF agents can be imported without any of
that installed, which is what lets CI test them.
"""

import logging
import os
import time
from concurrent.futures import ThreadPoolExecutor, as_completed

from .base_agent import AbstractStudyAgent
from .run_context import notes_system_block, strip_code_fence

log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Phase 2b — VideoAgent
# ---------------------------------------------------------------------------

class VideoAgent(AbstractStudyAgent):
    """Converts notes.md + section timing data into a narrated MP4 study video.

    Pipeline:
      1. _build_narration_scripts — expand timing entries into full narrations (LLM)
      2. _generate_html_slides    — LLM-generated HTML slides → Playwright PNG screenshots
      3. _generate_audio          — OpenAI TTS, one MP3 per narration
      4. _assemble_video          — MoviePy stitches frames + audio into MP4
    """

    def __init__(
        self,
        model: str = "claude-haiku-4-5-20251001",
        api_key: str | None = None,
        openai_api_key: str | None = None,
    ) -> None:
        """Initialise the agent.

        Args:
            model:          Anthropic model ID.
            api_key:        Anthropic API key.
            openai_api_key: OpenAI API key for TTS (falls back to OPENAI_API_KEY env var).
        """
        super().__init__(model=model, api_key=api_key)
        self.openai_api_key = openai_api_key

    def build_prompt(self, topic: str) -> str:
        """No-op — VideoAgent builds all prompts inline in its private methods.

        Args:
            topic: Unused; present to satisfy the abstract interface.

        Returns:
            An empty string.
        """
        return ""

    def run(
        self,
        notes_content: str,
        timing_json: list,
        output_dir: str | None = None,
    ) -> dict:
        """Generate a narrated study video and PPTX from notes and timing data.

        Args:
            notes_content: Full text of notes.md from Phase 1.
            timing_json:   List of section dicts, each with keys:
                           "section", "narration", "estimated_seconds".
            output_dir:    Absolute path to the run's output directory.
                           ``study_video.mp4`` is written directly here;
                           slides go into ``output_dir/slides/`` and audio
                           into ``output_dir/audio/``.
                           Falls back to ``output/`` at the project root when
                           None.

        Returns:
            ``{"status": "ok", "video_path": str, "pptx_path": str}``

        Raises:
            AssertionError: If frame count and narration count diverge after
                            slide generation.
        """
        if output_dir is None:
            _base = os.path.dirname(
                os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            )
            output_dir = os.path.join(_base, "output")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, "study_video.mp4")

        _t0 = time.monotonic()

        # Stage A: narrations and slides both need only timing_json — run in parallel
        t = time.monotonic()
        with ThreadPoolExecutor(max_workers=2) as _stage_a:
            _narr_fut  = _stage_a.submit(self._build_narration_scripts, notes_content, timing_json)
            _slide_fut = _stage_a.submit(self._generate_html_slides, timing_json, notes_content, output_dir)
            narrations  = _narr_fut.result()
            frame_paths = _slide_fut.result()
        # One wall-clock for the whole stage: narrations and slides ran in parallel.
        dt_narrations = time.monotonic() - t

        assert len(frame_paths) == len(narrations), (
            f"Frame/narration count mismatch: {len(frame_paths)} frames vs "
            f"{len(narrations)} narrations."
        )

        # Stage B+C: audio needs narrations; PPTX needs frame_paths — both ready, run in parallel
        t = time.monotonic()
        with ThreadPoolExecutor(max_workers=2) as _stage_bc:
            _audio_fut = _stage_bc.submit(self._generate_audio, narrations, output_dir)
            _pptx_fut  = _stage_bc.submit(self._export_pptx, frame_paths, timing_json, output_dir)
            audio_paths = _audio_fut.result()
            pptx_path   = _pptx_fut.result()
        # One wall-clock for the whole stage: audio and pptx ran in parallel.
        dt_audio = time.monotonic() - t

        # Stage D: video assembly needs both frame_paths and audio_paths
        t = time.monotonic()
        final_path = self._assemble_video(frame_paths, audio_paths, output_path)
        dt_assemble = time.monotonic() - t

        dt_total = time.monotonic() - _t0
        print(
            f"\n[Video] Benchmark ({len(timing_json)} sections):\n"
            f"  stage-A (narrations‖slides) : {dt_narrations:6.1f}s\n"
            f"  stage-BC (audio‖pptx)       : {dt_audio:6.1f}s\n"
            f"  assemble                    : {dt_assemble:6.1f}s\n"
            f"  total                       : {dt_total:6.1f}s"
        )

        result: dict = {
            "status": "ok",
            "video_path": final_path,
            "pptx_path": pptx_path,
            "benchmark": {
                "stage_a_s": round(dt_narrations, 2),   # narrations ‖ slides (wall-clock)
                "stage_bc_s": round(dt_audio, 2),        # audio ‖ pptx (wall-clock)
                "assemble_s": round(dt_assemble, 2),
                "total_s": round(dt_total, 2),
            },
        }
        if not self.validate_output(result):
            return {"status": "error", "video_path": final_path, "pptx_path": pptx_path}

        print(f"[Video] Saved to:      {final_path}")
        print(f"[Video] PPTX saved to: {pptx_path}")
        return result

    def validate_output(self, output: dict) -> bool:
        """Return True if both the MP4 and the PPTX exist and are non-empty.

        Args:
            output: The dict from run() — inspects ``output["video_path"]``
                    and ``output["pptx_path"]``.
        """
        video_path: str = output.get("video_path", "")
        pptx_path: str = output.get("pptx_path", "")
        video_ok = (
            isinstance(video_path, str)
            and video_path.endswith(".mp4")
            and os.path.isfile(video_path)
            and os.path.getsize(video_path) > 0
        )
        pptx_ok = (
            isinstance(pptx_path, str)
            and pptx_path.endswith(".pptx")
            and os.path.isfile(pptx_path)
            and os.path.getsize(pptx_path) > 0
        )
        return video_ok and pptx_ok

    # ------------------------------------------------------------------
    # Private helpers
    # ------------------------------------------------------------------

    def _build_narration_scripts(
        self, notes_content: str, timing_json: list
    ) -> list[str]:
        """Expand each timing entry into a full spoken narration via the LLM.

        All API calls are issued concurrently (one thread per section, capped
        at 5 workers to avoid rate-limit spikes).

        Args:
            notes_content: Full text of notes.md — context for every call.
            timing_json:   List of section dicts with keys "narration" and
                           "estimated_seconds".

        Returns:
            List of expanded narration strings, one per timing entry, in order.
        """
        # Built once and reused across every section. This loop previously
        # re-sent the whole notes text inside each per-section prompt, so a
        # ten-section video paid for them ten times. As a cached system block
        # the first call creates the entry and the rest read it.
        _system = notes_system_block(notes_content)

        def _narrate(i: int, entry: dict) -> tuple[int, str]:
            prompt = (
                f"Given this section from the notes: {entry['narration']}\n"
                f"Expand this into a spoken narration of approximately "
                f"{entry['estimated_seconds']} seconds.\n"
                "Tone: clear, educational, like a university lecturer.\n"
                "Return ONLY the narration text, no labels or headers."
            )
            return i, self._call_api(prompt, system=_system, use_tools=False)

        workers = min(len(timing_json), 5)
        narrations: list[str] = [""] * len(timing_json)
        with ThreadPoolExecutor(max_workers=workers) as pool:
            futures = {pool.submit(_narrate, i, entry): i for i, entry in enumerate(timing_json)}
            for fut in as_completed(futures):
                i, text = fut.result()
                narrations[i] = text
        return narrations

    def _generate_html_slides(
        self, timing_json: list, notes_content: str, output_dir: str
    ) -> list[str]:
        """Generate one HTML slide per timing entry and screenshot to PNG.

        LLM generation and Playwright screenshot for each slide are issued
        concurrently (one thread per slide, capped at 5 workers).

        Args:
            timing_json:   List of section dicts with keys "section" and
                           "narration".
            notes_content: Full notes text passed as LLM context.
            output_dir:    Absolute path to the run's output directory.
                           HTML and PNG files are placed in a ``slides/``
                           subdirectory within this directory.

        Returns:
            List of absolute PNG file paths, one per timing entry, in order.
        """
        # Byte-identical to the narration block, so slides read the cache entry
        # narration already created rather than making a second one. notes_content
        # was accepted here but never actually reached the model before.
        _system = notes_system_block(notes_content)

        def _make_slide(i: int, entry: dict) -> tuple[int, str]:
            prompt = (
                f"Generate a single self-contained HTML slide for the concept: "
                f"{entry['section']}\n"
                "The slide should:\n"
                "- Have a clean white background, 1280x720px\n"
                "- Show the section title in large text at top\n"
                "- Show 3-5 bullet points summarising the key ideas from the "
                f"narration: {entry['narration']}\n"
                "- Use simple inline CSS only — no external libraries\n"
                "- Be a complete HTML document that renders as a slide when "
                "opened in a browser\n"
                "Return ONLY the HTML, no explanation."
            )
            # "Return ONLY the HTML" is not binding, and an unstripped ```html
            # renders as body text before the doctype -- visible in the corner
            # of every slide, and so in every frame of the finished video.
            html_content = strip_code_fence(
                self._call_api(prompt, system=_system, use_tools=False)
            )
            html_path = self._save_output(
                html_content, f"slides/slide_{i:02d}.html", output_dir=output_dir
            )
            png_path = os.path.join(os.path.dirname(html_path), f"slide_{i:02d}.png")
            self._html_to_png(html_path, png_path)
            return i, png_path

        workers = min(len(timing_json), 5)
        png_paths: list[str] = [""] * len(timing_json)
        with ThreadPoolExecutor(max_workers=workers) as pool:
            futures = {pool.submit(_make_slide, i, entry): i for i, entry in enumerate(timing_json)}
            for fut in as_completed(futures):
                i, png_path = fut.result()
                png_paths[i] = png_path
        return png_paths

    def _html_to_png(self, html_path: str, output_png: str) -> str:
        """Convert an HTML file to a 1280×720 PNG using Playwright sync API.

        Uses the synchronous Playwright API so this method is safe to call
        from ThreadPoolExecutor threads (which have no running event loop)
        and from within ADK's async runner context (where asyncio.run would
        raise "This event loop is already running").

        Args:
            html_path:  Absolute path to the .html file.
            output_png: Destination path for the .png screenshot.

        Returns:
            output_png path after successful conversion.
        """
        from playwright.sync_api import sync_playwright

        with sync_playwright() as pw:
            browser = pw.chromium.launch()
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            page.goto(f"file://{html_path}")
            page.screenshot(path=output_png, full_page=False)
            browser.close()
        return output_png

    def _generate_audio(self, narrations: list[str], output_dir: str) -> list[str]:
        """Generate one TTS MP3 per narration using OpenAI tts-1-hd (voice: coral).

        All TTS requests are issued concurrently (one thread per narration,
        capped at 5 workers).

        Args:
            narrations: List of narration strings, one per slide.
            output_dir: Absolute path to the run's output directory.
                        MP3 files are placed in an ``audio/`` subdirectory
                        within this directory.

        Returns:
            List of MP3 file paths in the same order as narrations.
        """
        import openai

        audio_dir = os.path.join(output_dir, "audio")
        os.makedirs(audio_dir, exist_ok=True)

        client = openai.OpenAI(api_key=self.openai_api_key)

        def _tts(i: int, narration: str) -> tuple[int, str]:
            out_path = os.path.join(audio_dir, f"audio_{i:02d}.mp3")
            with client.audio.speech.with_streaming_response.create(
                model="tts-1-hd",
                voice="coral",
                input=narration,
            ) as response:
                response.stream_to_file(out_path)
            return i, out_path

        workers = min(len(narrations), 5)
        audio_paths: list[str] = [""] * len(narrations)
        with ThreadPoolExecutor(max_workers=workers) as pool:
            futures = {pool.submit(_tts, i, n): i for i, n in enumerate(narrations)}
            for fut in as_completed(futures):
                i, path = fut.result()
                audio_paths[i] = path
        return audio_paths

    def _export_pptx(
        self,
        frame_paths: list[str],
        timing_json: list,
        output_dir: str,
    ) -> str:
        """Build a widescreen PPTX where each slide is a PNG frame + title overlay.

        Each PNG frame is placed as a full-bleed background image. A dark navy
        text box is overlaid at the top of each slide with the section title in
        white 28 pt bold text.  The file is saved as
        ``{output_dir}/slides/flashcards.pptx``.

        Requires the ``python-pptx`` package (``pip install python-pptx``).

        Args:
            frame_paths: Absolute PNG file paths in slide order.
            timing_json: Section timing dicts; key ``"section"`` is used as
                         the slide title.
            output_dir:  Absolute path to the run's output directory.

        Returns:
            Absolute path to the saved ``flashcards.pptx``.
        """
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]  # truly blank layout

        for frame_path, entry in zip(frame_paths, timing_json):
            slide = prs.slides.add_slide(blank_layout)

            # Full-bleed PNG background
            slide.shapes.add_picture(
                frame_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

            # Title overlay — dark navy box at the top, white bold 28 pt text
            txBox = slide.shapes.add_textbox(
                left=0, top=0, width=prs.slide_width, height=Inches(1.2)
            )
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)  # dark navy
            tf = txBox.text_frame
            tf.word_wrap = True
            run = tf.paragraphs[0].add_run()
            run.text = entry.get("section", "")
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        slides_dir = os.path.join(output_dir, "slides")
        os.makedirs(slides_dir, exist_ok=True)
        pptx_path = os.path.join(slides_dir, "flashcards.pptx")
        prs.save(pptx_path)
        return pptx_path

    def _assemble_video(
        self,
        frame_paths: list[str],
        audio_paths: list[str],
        output_path: str,
    ) -> str:
        """Combine PNG frames and MP3 clips into a single MP4 via MoviePy.

        Each (frame, audio) pair becomes one clip whose on-screen duration
        matches the length of its audio clip.

        Args:
            frame_paths: PNG file paths, one per slide.
            audio_paths: MP3 file paths, one per slide.
            output_path: Destination MP4 path.

        Returns:
            output_path after successful export.
        """
        from moviepy import AudioFileClip, ImageClip, concatenate_videoclips

        clips = []
        for frame, audio in zip(frame_paths, audio_paths):
            audio_clip = AudioFileClip(audio)
            clips.append(
                ImageClip(frame)
                .with_duration(audio_clip.duration)
                .with_audio(audio_clip)
            )

        final = concatenate_videoclips(clips, method="compose")
        final.write_videofile(output_path, fps=24, codec="libx264", audio_codec="aac")
        return output_path
